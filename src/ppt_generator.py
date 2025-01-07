from pptx import Presentation
from datetime import datetime
from typing import List, Dict, Tuple
import os
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
from pptx.util import Pt

class PPTGeneratorError(Exception):
    pass

class PPTGenerator:
    def __init__(self):
        current_dir = os.path.dirname(os.path.abspath(__file__))
        self.template_path = os.path.join(current_dir, '..', 'resources', 'templates', 'template.pptx')
        
        if not os.path.exists(self.template_path):
            raise PPTGeneratorError(f"템플릿 파일을 찾을 수 없습니다: {self.template_path}")
            
        try:
            self.prs = Presentation(self.template_path)
            # 템플릿 구조 출력
            print(f"템플릿 슬라이드 수: {len(self.prs.slides)}")
            for idx, slide in enumerate(self.prs.slides):
                print(f"\n슬라이드 {idx + 1} 분석:")
                print(f"- 레이아웃: {slide.slide_layout.name}")
                print("- 도형 목록:")
                for shape_idx, shape in enumerate(slide.shapes):
                    print(f"  도형 {shape_idx + 1}:")
                    print(f"    유형: {shape.shape_type}")
                    if hasattr(shape, 'name'):
                        print(f"    이름: {shape.name}")
                    if shape.has_text_frame:
                        print(f"    텍스트: {shape.text}")
        except Exception as e:
            raise PPTGeneratorError(f"템플릿 분석 실패: {str(e)}")

    def create_birthday_slide(self, person: Dict) -> None:
        """생일자 슬라이드 생성"""
        try:
            template_slide = self.prs.slides[1]
            new_slide = self.prs.slides.add_slide(template_slide.slide_layout)
            
            # 템플릿의 도형들 복사
            for shape in template_slide.shapes:
                if hasattr(shape, 'shape_type'):
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        # 이미지 복사
                        image = new_slide.shapes.add_picture(
                            image_file=BytesIO(shape.image.blob),
                            left=left, top=top,
                            width=width, height=height
                        )
                    elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                        # 텍스트박스 복사
                        textbox = new_slide.shapes.add_textbox(
                            left=left, top=top,
                            width=width, height=height
                        )
                        
                        if shape.has_text_frame and shape.text:
                            # 원본 텍스트프레임과 새 텍스트프레임
                            orig_text_frame = shape.text_frame
                            new_text_frame = textbox.text_frame
                            
                            # 텍스트프레임 속성 복사
                            new_text_frame.word_wrap = orig_text_frame.word_wrap
                            
                            # 단락별로 복사
                            for i, orig_paragraph in enumerate(orig_text_frame.paragraphs):
                                if i == 0:
                                    new_paragraph = new_text_frame.paragraphs[0]
                                else:
                                    new_paragraph = new_text_frame.add_paragraph()
                                
                                # 단락 속성 복사
                                new_paragraph.alignment = orig_paragraph.alignment
                                new_paragraph.level = orig_paragraph.level
                                
                                # 텍스트 복사 및 치환
                                text = orig_paragraph.text
                                if text:
                                    birth_date = datetime.strptime(person['생년월일'], '%Y-%m-%d')
                                    text = text.replace("{name}", person['이름'])
                                    text = text.replace("{month}", str(birth_date.month))
                                    text = text.replace("{day}", str(birth_date.day))
                                    
                                    new_paragraph.text = text
                                    
                                    # 런(서식 단위)별로 복사
                                    if len(orig_paragraph.runs) > 0:
                                        orig_run = orig_paragraph.runs[0]
                                        new_run = new_paragraph.runs[0]
                                        
                                        # 폰트 속성 복사
                                        new_font = new_run.font
                                        orig_font = orig_run.font
                                        
                                        # 기본 폰트 설정
                                        new_font.name = "Maplestory OTF"
                                        
                                        # 크기 복사
                                        if orig_font.size is not None:
                                            new_font.size = orig_font.size
                                        
                                        # 색상 복사 (테마 색상 처리)
                                        try:
                                            if hasattr(orig_font.color, 'rgb'):
                                                new_font.color.rgb = orig_font.color.rgb
                                            elif hasattr(orig_font.color, 'theme_color'):
                                                new_font.color.theme_color = orig_font.color.theme_color
                                        except Exception as color_error:
                                            print(f"색상 복사 중 오류 (무시됨): {str(color_error)}")
                                        
                                        # 기타 폰트 속성 복사
                                        new_font.bold = orig_font.bold
                                        new_font.italic = orig_font.italic
                                        new_font.underline = orig_font.underline
                                        
                                        print(f"폰트 정보 - {text}:")
                                        print(f"- 이름: {new_font.name}")
                                        print(f"- 크기: {new_font.size}")
                                        if hasattr(new_font.color, 'theme_color'):
                                            print(f"- 색상: Theme Color")
                                        elif hasattr(new_font.color, 'rgb'):
                                            print(f"- 색상: RGB")
                            
            print(f"{person['이름']}의 슬라이드 생성 완료")
                            
        except Exception as e:
            print(f"슬라이드 생성 중 오류: {str(e)}")
            raise PPTGeneratorError(f"슬라이드 생성 오류: {str(e)}")
        
    def create_title_slide(self, month: int) -> None:
        """월별 타이틀 슬라이드 수정"""
        try:
            print(f"\n타이틀 슬라이드 수정 (월: {month})")
            title_slide = self.prs.slides[0]
            
            print("현재 도형들의 텍스트:")
            for shape in title_slide.shapes:
                if shape.has_text_frame:
                    print(f"- {shape.text}")
                        
            for shape in title_slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    
                    # 단락별로 처리
                    for i, paragraph in enumerate(text_frame.paragraphs):
                        # 원본 속성 저장
                        original_font = None
                        if len(paragraph.runs) > 0:
                            original_font = paragraph.runs[0].font

                        if "{month}" in paragraph.text:
                            original_text = paragraph.text
                            new_text = original_text.replace("{month}", str(month))
                            
                            # 텍스트 설정
                            paragraph.text = new_text
                            
                            # 새로운 런에 원본 속성 적용
                            if len(paragraph.runs) > 0 and original_font:
                                new_run = paragraph.runs[0]
                                new_font = new_run.font
                                
                                # 기본 폰트 설정
                                new_font.name = "Maplestory OTF"
                                
                                # 크기 복사
                                if original_font.size is not None:
                                    new_font.size = original_font.size
                                
                                # 색상 복사
                                try:
                                    if hasattr(original_font.color, 'rgb'):
                                        new_font.color.rgb = original_font.color.rgb
                                    elif hasattr(original_font.color, 'theme_color'):
                                        new_font.color.theme_color = original_font.color.theme_color
                                except Exception as color_error:
                                    print(f"타이틀 색상 복사 중 오류 (무시됨): {str(color_error)}")
                                
                                # 기타 속성 복사
                                new_font.bold = original_font.bold
                                new_font.italic = original_font.italic
                                new_font.underline = original_font.underline
                                
                                print(f"타이틀 폰트 정보 - {new_text}:")
                                print(f"- 이름: {new_font.name}")
                                print(f"- 크기: {new_font.size}")
                                if hasattr(new_font.color, 'theme_color'):
                                    print(f"- 색상: Theme Color")
                                elif hasattr(new_font.color, 'rgb'):
                                    print(f"- 색상: RGB")
                            
                            print(f"텍스트 교체: {original_text} -> {new_text}")
                        else:
                            # month가 포함되지 않은 텍스트(HAPPY BIRTHDAY 등)도 폰트 적용
                            if len(paragraph.runs) > 0:
                                run = paragraph.runs[0]
                                font = run.font
                                font.name = "Maplestory OTF"
                                
                                if original_font:
                                    # 원본 속성 복사
                                    if original_font.size is not None:
                                        font.size = original_font.size
                                    try:
                                        if hasattr(original_font.color, 'rgb'):
                                            font.color.rgb = original_font.color.rgb
                                        elif hasattr(original_font.color, 'theme_color'):
                                            font.color.theme_color = original_font.color.theme_color
                                    except Exception:
                                        pass
                                    font.bold = original_font.bold
                                    font.italic = original_font.italic
                                    font.underline = original_font.underline

        except Exception as e:
            print(f"타이틀 슬라이드 수정 중 오류: {str(e)}")
            raise PPTGeneratorError(f"타이틀 슬라이드 수정 오류: {str(e)}")

    def generate_ppt(self, month: int, birthday_list: List[Dict], save_path: str) -> Tuple[bool, str]:
        try:
            print(f"\nPPT 생성 시작:")
            print(f"- 월: {month}")
            print(f"- 생일자 수: {len(birthday_list)}")
            print(f"- 저장 경로: {save_path}")
            
            self._validate_save_path(save_path)
            self._validate_birthday_data(birthday_list)
            
            self.create_title_slide(month)
            
            for person in birthday_list:
                self.create_birthday_slide(person)
            
            # 템플릿 슬라이드 제거
            xml_slides = self.prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[1])
            print("템플릿 슬라이드 제거됨")
            
            output_path = os.path.join(save_path, f"{month}월_생일자.pptx")
            self.prs.save(output_path)
            print(f"파일 저장 완료: {output_path}")
            
            return True, f"PPT 파일이 생성되었습니다: {output_path}"
            
        except Exception as e:
            print(f"PPT 생성 실패: {str(e)}")
            return False, f"PPT 생성 실패: {str(e)}"

    def _validate_save_path(self, save_path: str) -> None:
        if not os.path.exists(save_path):
            raise PPTGeneratorError(f"저장 경로가 존재하지 않습니다: {save_path}")
        if not os.access(save_path, os.W_OK):
            raise PPTGeneratorError(f"저장 경로에 쓰기 권한이 없습니다: {save_path}")

    def _validate_birthday_data(self, birthday_list: List[Dict]) -> None:
        if not birthday_list:
            raise PPTGeneratorError("생일자 데이터가 비어있습니다")
        
        required_fields = {'이름', '성별', '생년월일', '나이'}
        for person in birthday_list:
            missing_fields = required_fields - set(person.keys())
            if missing_fields:
                raise PPTGeneratorError(f"필수 필드가 누락되었습니다: {', '.join(missing_fields)}")